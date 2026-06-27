"""Generate the pitch deck for the North Range Mine gas-safety project.

Theme: simple, readable, mining-flavoured — deep slate navy + safety amber on
white. Run:  backend/venv/Scripts/python.exe make_deck.py
Output:  North_Range_Mine_Safety.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ---------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x41)   # deep slate navy (primary)
SLATE = RGBColor(0x2E, 0x40, 0x57)  # secondary
AMBER = RGBColor(0xF4, 0xA3, 0x00)  # safety amber (accent)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)  # off-white background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1B, 0x2A, 0x41)    # body text
GREY = RGBColor(0x5B, 0x66, 0x78)   # muted text
GREEN = RGBColor(0x2E, 0x7D, 0x32)
YELLOW = RGBColor(0xF4, 0xA3, 0x00)
RED = RGBColor(0xC6, 0x28, 0x2E)

FONT = "Segoe UI"
FONT_L = "Segoe UI Light"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---- low-level helpers -----------------------------------------------------
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_autofit(tf):
    # keep text from auto-resizing the box weirdly
    tf.word_wrap = True


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_fill(s, color)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (str, size, bold, color, font)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, bold, color, font) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=18, color=INK, gap=10, bullet_color=AMBER):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    _no_autofit(tf)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        # amber square bullet
        rb = p.add_run()
        rb.text = "\u25A0  "
        rb.font.size = Pt(size - 4)
        rb.font.color.rgb = bullet_color
        rb.font.name = FONT
        # support (head, tail) tuples for a bold lead-in
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run(); r1.text = head
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = color; r1.font.name = FONT
            r2 = p.add_run(); r2.text = tail
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = FONT
        else:
            r = p.add_run(); r.text = item
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def header(slide, title, kicker=None):
    """Standard content-slide header: navy bar + amber underline + title."""
    rect(slide, 0, 0, 13.333, 1.15, NAVY)
    rect(slide, 0, 1.15, 13.333, 0.06, AMBER)
    # small amber tick on the left
    rect(slide, 0.55, 0.34, 0.12, 0.5, AMBER)
    if kicker:
        text(slide, 0.85, 0.20, 11.5, 0.35,
             [[(kicker.upper(), 12, True, AMBER, FONT)]])
        text(slide, 0.85, 0.46, 12.0, 0.6,
             [[(title, 26, True, WHITE, FONT)]])
    else:
        text(slide, 0.85, 0.30, 12.0, 0.7,
             [[(title, 28, True, WHITE, FONT)]], anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, n):
    text(slide, 0.55, 7.0, 9.0, 0.4,
         [[("North Range Mine \u2014 Offline Edge Gas-Safety", 10, False, GREY, FONT)]])
    text(slide, 11.8, 7.0, 1.0, 0.4,
         [[(f"{n:02d}", 10, True, AMBER, FONT)]], align=PP_ALIGN.RIGHT)


def bg(slide, color=LIGHT):
    rect(slide, -0.05, -0.05, 13.45, 7.6, color)


def chip(slide, x, y, w, label, color):
    s = rect(slide, x, y, w, 0.55, color, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    return s


def flow_box(slide, x, y, w, h, title, sub, fill=WHITE, tcolor=INK, border=AMBER):
    s = rect(slide, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE)
    s.line.color.rgb = border
    s.line.width = Pt(1.5)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = tcolor; r.font.name = FONT
    for line in sub.split("\n"):
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = line
        r2.font.size = Pt(10.5); r2.font.color.rgb = GREY; r2.font.name = FONT
    return s


def arrow(slide, x, y, w, h=0.4, color=AMBER):
    s = rect(slide, x, y, w, h, color, MSO_SHAPE.RIGHT_ARROW)
    return s


def card(slide, x, y, w, h, fill=WHITE, border=RGBColor(0xE1, 0xE6, 0xEC)):
    s = rect(slide, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE)
    s.line.color.rgb = border
    s.line.width = Pt(1.0)
    return s


# ===========================================================================
# SLIDE 1 — TITLE
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, 0.35, 7.5, AMBER)                 # left spine
rect(s, 0.9, 2.05, 1.6, 0.14, AMBER)            # accent line
text(s, 0.9, 1.35, 11, 0.5, [[("UNDERGROUND MINE SAFETY  \u2022  EDGE  \u2022  OFFLINE", 14, True, AMBER, FONT)]])
text(s, 0.86, 2.25, 11.6, 2.2, [
    [("North Range Mine", 52, True, WHITE, FONT)],
    [("Offline edge gas-safety early-warning", 32, False, RGBColor(0xCF, 0xD8, 0xE3), FONT_L)],
])
text(s, 0.9, 4.7, 11.6, 1.2, [
    [("Turns a gas reading into an immediate, cited safety action \u2014 per level,", 18, False, RGBColor(0xB9, 0xC4, 0xD2), FONT)],
    [("at the edge, with no internet and no central server to fail.", 18, False, RGBColor(0xB9, 0xC4, 0xD2), FONT)],
])
chip(s, 0.9, 6.25, 3.2, "Track: Mining / Sustainability", SLATE)
chip(s, 4.3, 6.25, 2.3, "3-hour build", SLATE)
chip(s, 6.8, 6.25, 3.0, "Runs with Wi-Fi OFF", AMBER)

# ===========================================================================
# SLIDE 2 — THE PROBLEM (criterion 1: problem relevance + track)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "No signal. No second chances.", kicker="The Problem")
bullets(s, 0.85, 1.55, 7.4, 4.8, [
    ("Underground has no internet, no cellular. ", "Rock blocks radio \u2014 it's the defining constraint, not an edge case."),
    ("Leaky-feeder cable dies on a roof fall. ", "Any system that leans on a central server goes dark exactly during an incident."),
    ("Sensors report a number, not an action. ", "\u201cMethane 1.4%\u201d doesn't tell the crew what to do, or cite why."),
    ("Alarms fire too late. ", "A fixed threshold reacts after the danger line is crossed \u2014 not while it's climbing."),
], size=18, gap=16)
# right stat card
card(s, 8.7, 1.7, 3.9, 4.4)
rect(s, 8.7, 1.7, 3.9, 0.12, AMBER)
text(s, 9.0, 2.05, 3.4, 4.0, [
    [("Why it matters", 16, True, NAVY, FONT)],
    [("", 6, False, GREY, FONT)],
    [("Gas builds in seconds", 22, True, RED, FONT)],
    [("after a blast or strata release.", 13, False, GREY, FONT)],
    [("", 8, False, GREY, FONT)],
    [("Comms fail first", 22, True, AMBER, FONT)],
    [("in the moment you rely on them.", 13, False, GREY, FONT)],
    [("", 8, False, GREY, FONT)],
    [("The crew needs the answer", 18, True, GREEN, FONT)],
    [("at the rock face, offline.", 13, False, GREY, FONT)],
])
footer(s, 2)

# ===========================================================================
# SLIDE 3 — OUR SOLUTION (criterion 2: solution design)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "An early-warning brain at the rock face", kicker="Our Solution")
text(s, 0.85, 1.5, 11.6, 0.9, [
    [("Runs on a small edge box on each level \u2014 fully offline. It reads the air, sees the ", 18, False, INK, FONT),
     ("trend", 18, True, NAVY, FONT),
     (", and ", 18, False, INK, FONT),
     ("speaks the exact regulation", 18, True, NAVY, FONT),
     (".", 18, False, INK, FONT)],
])
# four pillars
px, py, pw, gap = 0.85, 2.6, 2.78, 0.18
pillars = [
    ("Per level", "Each level monitored on its own \u2014 no single point of failure.", NAVY),
    ("Trend, not threshold", "Flags air climbing toward danger so crews leave early.", AMBER),
    ("Cited action", "Speaks the exact O. Reg 854 step, with the citation.", GREEN),
    ("Fully offline", "MQTT mesh + edge brain + voice. Wi-Fi can be off.", SLATE),
]
for i, (t, d, c) in enumerate(pillars):
    x = px + i * (pw + gap)
    card(s, x, py, pw, 3.1)
    rect(s, x, py, pw, 0.5, c)
    text(s, x + 0.05, py + 0.06, pw - 0.1, 0.4, [[(t, 15, True, WHITE, FONT)]], align=PP_ALIGN.CENTER)
    text(s, x + 0.22, py + 0.75, pw - 0.44, 2.2, [[(d, 13.5, False, INK, FONT)]])
footer(s, 3)

# ===========================================================================
# SLIDE 4 — HOW IT WORKS (criterion 2/3): architecture + data flow
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "How it works \u2014 data flow, all on one box", kicker="Architecture")
# flow boxes
y = 2.05
bw, bh = 2.25, 1.5
xs = [0.7, 3.35, 6.0, 8.65, 11.3 - 2.25 + 0.0]
flow_box(s, 0.7, y, bw, bh, "Level sensors", "CH4 / CO / NO2 / O2\nairflow \u2022 fan", border=NAVY)
arrow(s, 3.02, y + 0.55, 0.5)
flow_box(s, 3.55, y, bw, bh, "Local mesh", "MQTT broker\n(or internal timer)", border=AMBER)
arrow(s, 5.87, y + 0.55, 0.5)
flow_box(s, 6.4, y, bw, bh, "Edge brain", "classify + trend\nper level", border=GREEN)
arrow(s, 8.72, y + 0.55, 0.5)
flow_box(s, 9.25, y, bw, bh, "Action + voice", "cited Reg 854\nspoken aloud", border=RED)
# historian below
card(s, 0.7, 4.05, 11.8, 1.0, RGBColor(0xEE, 0xF2, 0xF7))
rect(s, 0.7, 4.05, 0.12, 1.0, AMBER)
text(s, 1.0, 4.18, 11.3, 0.8, [
    [("Offline historian \u2014 ", 15, True, NAVY, FONT),
     ("every reading + event appended to dated files on disk (telemetry / events / app log). Survives restarts; never leaves the machine.", 14, False, INK, FONT)],
])
text(s, 0.7, 5.35, 11.8, 1.2, [
    [("Two run modes, identical output: ", 14, True, NAVY, FONT),
     ("MQTT mesh (the real story) or a broker-free internal timer (the demo safety net). The frontend polls every 2\u20133 s over localhost.", 14, False, INK, FONT)],
])
text(s, 0.7, 6.05, 11.8, 0.6, [
    [("Every arrow is on 127.0.0.1 \u2014 loopback only. Pull the network cable and nothing changes.", 13, True, GREEN, FONT)],
])
footer(s, 4)

# ===========================================================================
# SLIDE 5 — BUILT LIKE A REAL MINE (criterion 3 + 7)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Modelled on a real mine", kicker="Domain Realism")
bullets(s, 0.85, 1.55, 7.3, 5.0, [
    ("Levels, not zones. ", "Horizontal horizons named by depth (400L\u20131600L) on an intake\u2013working\u2013return air circuit."),
    ("Multi-gas station. ", "CH4, CO, CO2, NO2, O2 (low = hazard) \u2014 worst gas drives the status, like a real fixed monitor."),
    ("Drill-and-blast cycle. ", "A blast spikes CO/NO2; ventilation-on-demand ramps the fan and clears it."),
    ("CO governs re-entry. ", "CO clears slowest, so re-entry waits on CO \u2014 exactly as in practice."),
    ("Rolling CO average. ", "A TWA-style figure backs the re-entry decision."),
], size=17, gap=14)
# status ladder
card(s, 8.6, 1.7, 4.0, 4.5)
text(s, 8.85, 1.85, 3.6, 0.5, [[("Active level: 1200L", 16, True, NAVY, FONT)]])
text(s, 8.85, 2.3, 3.6, 0.4, [[("Live drill-and-blast loop", 12, False, GREY, FONT)]])
chip(s, 8.85, 2.85, 3.5, "GREEN \u2014 crew working, air clear", GREEN)
chip(s, 8.85, 3.55, 3.5, "RED \u2014 blast: CO ~220 ppm", RED)
chip(s, 8.85, 4.25, 3.5, "YELLOW \u2014 clearing (fan 100%)", YELLOW)
chip(s, 8.85, 4.95, 3.5, "GREEN \u2014 re-entry granted", GREEN)
text(s, 8.85, 5.7, 3.6, 0.5, [[("\u2193 loops continuously for the demo", 12, True, AMBER, FONT)]])
footer(s, 5)

# ===========================================================================
# SLIDE 6 — WORKING PROTOTYPE (criterion 3: technical execution)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "What actually runs today", kicker="Working Prototype")
# left: stack
text(s, 0.85, 1.5, 6.0, 0.5, [[("Tech stack \u2014 all offline-capable", 16, True, NAVY, FONT)]])
bullets(s, 0.85, 2.0, 6.1, 4.4, [
    ("FastAPI edge backend ", "\u2014 /levels, /alert, /ask, /health (+ /zones alias)."),
    ("MQTT (Mosquitto) ", "as the mine mesh; paho publisher per level."),
    ("Pydantic-validated readings ", "\u2014 one bad payload can't crash monitoring."),
    ("Template Reg 854 answers ", "\u2014 no LLM: instant, deterministic, never hallucinates."),
    ("Offline keyword retrieval ", "over local docs for citations."),
    ("Browser Web Speech ", "for the spoken alert \u2014 also offline."),
], size=15, gap=11)
# right: demo arc
card(s, 7.3, 1.5, 5.3, 4.9)
rect(s, 7.3, 1.5, 5.3, 0.55, NAVY)
text(s, 7.5, 1.6, 5.0, 0.4, [[("60-second demo arc", 16, True, WHITE, FONT)]])
steps = [
    "Level map all green",
    "1200L blasts \u2192 CO spikes, turns RED",
    "Banner fires; voice speaks the cited CO procedure",
    "Fan ramps to 100% \u2014 air clears (YELLOW)",
    "CO back under limit \u2192 re-entry granted (GREEN)",
    "Turn Wi-Fi OFF \u2014 do it again, still works",
]
ty = 2.25
for i, st in enumerate(steps):
    rect(s, 7.55, ty + 0.04, 0.32, 0.32, AMBER, MSO_SHAPE.OVAL)
    text(s, 7.55, ty + 0.0, 0.32, 0.32, [[(str(i + 1), 12, True, NAVY, FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 8.0, ty - 0.02, 4.5, 0.6, [[(st, 13.5, False, INK, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    ty += 0.68
footer(s, 6)

# ===========================================================================
# SLIDE 7 — IMPACT & FEASIBILITY (criterion 4)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Impact & feasibility", kicker="Why it's real")
cols = [
    ("Saves lives", "Orderly withdrawal on a rising trend \u2014 before an emergency evacuation.", GREEN),
    ("Saves energy", "Ventilation-on-demand runs fans only when needed (vent can be >50% of mine energy).", AMBER),
    ("Deployable", "Same code on a laptop or an underground edge box \u2014 config only, no rewrite.", NAVY),
]
for i, (t, d, c) in enumerate(cols):
    x = 0.85 + i * 4.0
    card(s, x, 1.7, 3.7, 2.5)
    rect(s, x, 1.7, 3.7, 0.5, c)
    text(s, x + 0.05, 1.77, 3.6, 0.4, [[(t, 16, True, WHITE, FONT)]], align=PP_ALIGN.CENTER)
    text(s, x + 0.25, 2.45, 3.2, 1.6, [[(d, 14, False, INK, FONT)]])
text(s, 0.85, 4.5, 11.6, 0.5, [[("Next steps", 16, True, NAVY, FONT)]])
bullets(s, 0.85, 5.0, 11.6, 1.8, [
    ("Real sensor nodes ", "on a self-healing wireless mesh between levels."),
    ("On-device TTS ", "in the cap lamp / tablet; pilot on one active level."),
    ("Sync at surface ", "\u2014 the historian uploads only when a backhaul link is available."),
], size=15, gap=9)
footer(s, 7)

# ===========================================================================
# SLIDE 8 — RESPONSIBLE DESIGN & LIMITATIONS (criterion 6)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Responsible design & honest limits", kicker="Awareness")
text(s, 0.85, 1.5, 5.7, 0.5, [[("Designed responsibly", 16, True, GREEN, FONT)]])
bullets(s, 0.85, 2.0, 5.8, 4.6, [
    ("Privacy by design ", "\u2014 data stays underground; no cloud, no tracking."),
    ("Safety-first answers ", "\u2014 template + citation, not a guessing model."),
    ("Fail-safe ", "\u2014 keeps serving last state if the broker drops."),
    ("Auditable ", "\u2014 every reading + decision logged offline."),
], size=15, gap=11)
text(s, 6.95, 1.5, 5.6, 0.5, [[("Limitations we own", 16, True, RED, FONT)]])
bullets(s, 6.95, 2.0, 5.6, 4.6, [
    ("Demo thresholds ", "are tuned to be watchable \u2014 not legal limits."),
    ("Sensors are simulated; ", "real calibration/drift handling is future work."),
    ("Polling, not push ", "\u2014 2\u20133 s cadence; fine for demo, not hard real-time."),
    ("No real radio mesh yet ", "\u2014 MQTT stands in for it on one machine."),
], size=15, gap=11, bullet_color=RED)
footer(s, 8)

# ===========================================================================
# SLIDE 9 — ORIGINALITY / WHAT MAKES US DIFFERENT (criterion 7)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "What makes us different", kicker="Originality")
text(s, 0.85, 1.45, 11.6, 0.8, [
    [("Established vendors (Maestro, Newtrax, Becker) ", 16, True, NAVY, FONT),
     ("sense gas and show the number. We sit on top and turn it into a cited action.", 16, False, INK, FONT)],
])
# comparison two cards
card(s, 0.85, 2.35, 5.7, 3.9)
rect(s, 0.85, 2.35, 5.7, 0.55, SLATE)
text(s, 1.05, 2.45, 5.3, 0.4, [[("Typical systems", 15, True, WHITE, FONT)]])
bullets(s, 1.1, 3.1, 5.2, 3.0, [
    "Display the gas value on a dashboard",
    "Alarm after a fixed limit is crossed",
    "Lean on a surface server / control room",
], size=14, gap=12, bullet_color=GREY)
card(s, 6.95, 2.35, 5.65, 3.9)
rect(s, 6.95, 2.35, 5.65, 0.55, AMBER)
text(s, 7.15, 2.45, 5.3, 0.4, [[("Ours \u2014 unique combination", 15, True, NAVY, FONT)]])
bullets(s, 7.2, 3.1, 5.2, 3.0, [
    "Edge-first & offline, per level (no central failure point)",
    "Trend-based early warning (leave before the line)",
    "Speaks the exact cited Reg 854 procedure",
    "Offline append-only audit historian",
], size=14, gap=10, bullet_color=GREEN)
footer(s, 9)

# ===========================================================================
# SLIDE 10 — CLOSING / PITCH LINE (criterion 5)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, 0.35, 7.5, AMBER)
rect(s, 0.9, 2.0, 1.6, 0.14, AMBER)
text(s, 0.9, 1.3, 11, 0.5, [[("THE PITCH", 14, True, AMBER, FONT)]])
text(s, 0.86, 2.25, 11.6, 3.2, [
    [("\u201cThe moment a level trends toward dangerous \u2014", 26, False, WHITE, FONT_L)],
    [("not just when it crosses the line \u2014 it speaks the exact", 26, False, WHITE, FONT_L)],
    [("regulation, offline, at the rock face.", 26, False, WHITE, FONT_L)],
    [("The crew gets out before it's an emergency.\u201d", 26, True, AMBER, FONT)],
])
chip(s, 0.9, 6.0, 3.6, "Demo: Wi-Fi OFF, it still works", AMBER)
chip(s, 4.7, 6.0, 2.6, "Thank you", SLATE)
footer(s, 10)

# ===========================================================================
# SLIDE 11 — APPENDIX: criteria coverage map
# ===========================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "How this maps to the judging criteria", kicker="Appendix")
rows = [
    ("1  Problem relevance & track", "Mining safety, the no-internet constraint (slide 2)"),
    ("2  Solution design & practicality", "Edge brain: per-level, trend, cited action (slides 3\u20134)"),
    ("3  Working prototype & execution", "FastAPI + MQTT + voice, live demo (slides 4\u20136)"),
    ("4  Impact & feasibility", "Lives, energy, edge-deployable, next steps (slide 7)"),
    ("5  Presentation & pitch", "60-sec arc + the pitch line (slides 6, 10)"),
    ("6  Responsible design & limits", "Privacy, fail-safe, honest limitations (slide 8)"),
    ("7  Originality & effort", "Cited-action edge angle vs vendors (slides 5, 9)"),
]
ty = 1.55
for i, (crit, where) in enumerate(rows):
    fill = WHITE if i % 2 == 0 else RGBColor(0xEE, 0xF2, 0xF7)
    card(s, 0.85, ty, 11.65, 0.68, fill, fill)
    rect(s, 0.85, ty, 0.1, 0.68, AMBER)
    text(s, 1.1, ty + 0.06, 5.0, 0.56, [[(crit, 14, True, NAVY, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.2, ty + 0.06, 6.2, 0.56, [[(where, 13, False, INK, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    ty += 0.74
footer(s, 11)

out = "North_Range_Mine_Safety.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
