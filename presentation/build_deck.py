"""Generate hackathon pitch deck aligned to official judging criteria."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --- Brand palette (mining / safety) ---
BG = RGBColor(15, 23, 42)
ACCENT = RGBColor(245, 158, 11)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 68)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(148, 163, 184)
BLUE = RGBColor(56, 189, 248)

OUTPUT = Path(__file__).parent / "MineGuard_Edge_Hackathon.pptx"


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=WHITE):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_accent_bar(slide, top=0.55):
    shape = slide.shapes.add_shape(1, Inches(0.6), Inches(top), Inches(0.08), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def slide_header(slide, title, subtitle=None, criterion=None):
    add_accent_bar(slide)
    if criterion:
        add_textbox(slide, 0.85, 0.35, 11.5, 0.35, criterion, size=11, bold=True, color=BLUE)
        title_top = 0.65
        sub_top = 1.25
    else:
        title_top = 0.45
        sub_top = 1.05
    add_textbox(slide, 0.85, title_top, 11.5, 0.7, title, size=32, bold=True)
    if subtitle:
        add_textbox(slide, 0.85, sub_top, 11.5, 0.5, subtitle, size=14, color=MUTED)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── TITLE ──────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.9, 1.6, 11.5, 1.2, "MineGuard Edge", size=48, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        s, 0.9, 2.7, 11.5, 0.8,
        "Offline Gas Safety Intelligence for Underground Mines",
        size=22, color=ACCENT, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, 0.9, 3.6, 11.5, 0.5,
        "Hackathon Track:  MINING  ·  also touches Health (worker safety) & Sustainability (ventilation efficiency)",
        size=14, color=GREEN, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, 0.9, 4.5, 11.5, 0.6,
        "Team of 3  ·  100% Local / No Cloud  ·  github.com/RavalDH/Cursor_Hack",
        size=13, color=MUTED, align=PP_ALIGN.CENTER,
    )

    # ── CRITERION 1: Problem relevance & track alignment ─────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(
        s,
        "The Problem",
        "Specific, real, and directly aligned to the Mining track",
        criterion="JUDGING CRITERION 1 — Problem relevance & track alignment",
    )
    add_bullets(s, 0.9, 1.85, 11.2, 3.8, [
        "Underground mines have NO internet and NO cellular — rock blocks every signal",
        "Sensors report a gas number — they do not tell crews what to do next",
        "Cloud systems go dark during roof falls, exactly when they are needed most",
        "Alarms fire AFTER danger — not while methane is still climbing toward the limit",
    ], size=19)
    add_textbox(s, 0.9, 5.0, 11.5, 0.4, "Track alignment", size=16, bold=True, color=ACCENT)
    add_bullets(s, 0.9, 5.4, 11.2, 1.5, [
        "MINING — gas monitoring, ventilation, Reg 854 compliance underground",
        "HEALTH — protects miners with early warning and cited evacuation procedures",
        "SUSTAINABILITY — smart fan control saves power; only ramps when gas builds",
    ], size=15, color=MUTED)

    # ── CRITERION 2: Solution design & practicality ────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(
        s,
        "Our Solution",
        "Designed for real mine operators, supervisors, and underground crews",
        criterion="JUDGING CRITERION 2 — Solution design & practicality",
    )
    add_textbox(
        s, 0.9, 1.85, 11.5, 1.0,
        "MineGuard Edge runs offline at every zone, detects dangerous gas trends early, "
        "ramps ventilation automatically, and speaks the exact Reg 854 procedure — with citations.",
        size=21,
    )
    add_bullets(s, 0.9, 3.1, 5.5, 3.5, [
        "Who uses it: shift supervisors, safety officers, underground crews",
        "Where it runs: edge box per zone — no surface server required",
        "When it acts: trend detected → fan ramps → alert + voice procedure",
        "Why it fits: same MQTT + mesh architecture Becker / Maestro / Glencore use",
    ], size=16)
    add_textbox(s, 7.0, 3.1, 5.5, 3.5,
        "Practical design choices:\n\n"
        "• Trend before threshold → orderly withdrawal\n"
        "• Ventilation only (never enrich O₂ — explosion risk)\n"
        "• Template answers over cloud LLM → instant, reliable\n"
        "• Internal-timer fallback if MQTT is unavailable",
        size=15, color=ACCENT,
    )

    # ── Architecture (supports criterion 2 & 3) ──────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(s, "How It Works", "All local — no external calls anywhere")
    arch = (
        "ZONE SENSORS (×4)          EDGE BACKEND\n"
        "  methane, CO, temp   →    Mosquitto MQTT (mesh stand-in)\n"
        "  airflow, fan speed         ↓\n"
        "                             Trend detector + fan control loop\n"
        "                             ↓\n"
        "                             FastAPI  /zones  /alert  /ask\n"
        "                             ↓\n"
        "                             Reg 854 docs + cited answers\n"
        "                             ↓\n"
        "                             Web UI + offline voice alert"
    )
    box = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(7.2), Inches(4.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = arch
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    add_bullets(s, 8.4, 1.9, 4.2, 4.5, [
        "Production: wireless self-healing mesh",
        "Demo: simulated zones on one laptop",
        "Green → Yellow → Red → ventilating → Green cycle",
    ], size=15, color=MUTED)

    # ── Key differentiator ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(s, "Key Innovation: TREND, Not Threshold", "Early warning — crew leaves before it's an emergency")
    add_textbox(s, 0.9, 1.7, 5.5, 0.4, "Traditional systems", size=17, bold=True, color=RED)
    add_bullets(s, 0.9, 2.15, 5.5, 1.8, ["Alarm only at 1.5% methane", "Reactive — too late for orderly exit"], size=16)
    add_textbox(s, 7.0, 1.7, 5.5, 0.4, "MineGuard Edge", size=17, bold=True, color=GREEN)
    add_bullets(s, 7.0, 2.15, 5.5, 2.2, [
        "Track last N readings per zone",
        "Rising at 0.9%+ → YELLOW warning",
        "At 1.5%+ → RED + Reg 854 evacuation cited aloud",
    ], size=16)
    add_textbox(
        s, 0.9, 5.0, 11.5, 1.5,
        "Ventilation loop: fan ramps to dilute gas → zone clears → fan eases back to idle → all-clear voice message",
        size=16, color=ACCENT, align=PP_ALIGN.CENTER,
    )

    # ── CRITERION 3: Working prototype & technical execution ─────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(
        s,
        "Working Prototype",
        "Live, runnable, and demonstrable right now",
        criterion="JUDGING CRITERION 3 — Working prototype & technical execution",
    )
    add_bullets(s, 0.9, 1.85, 5.8, 4.5, [
        "FastAPI backend — running at localhost:8000",
        "4 simulated zones with live gas + fan telemetry",
        "Trend detector + ventilation control loop implemented",
        "Reg 854 keyword retrieval + /ask endpoint",
        "Recovery detection — all-clear when zone returns to green",
    ], size=16)
    add_textbox(s, 7.0, 1.85, 5.5, 0.4, "API endpoints (live)", size=16, bold=True, color=ACCENT)
    endpoints = "GET  /zones   → zone map data\nGET  /alert   → cited safety alert\nPOST /ask     → Reg 854 Q&A\nGET  /health  → system status"
    box = s.shapes.add_textbox(Inches(7.0), Inches(2.35), Inches(5.5), Inches(2.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = endpoints
    p.font.name = "Courier New"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    add_textbox(
        s, 0.9, 5.8, 11.5, 0.8,
        "Tech: Python · FastAPI · Mosquitto MQTT · paho-mqtt · Web Speech API · Lovable UI",
        size=14, color=MUTED,
    )

    # ── CRITERION 5: Presentation & pitch — demo arc ───────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(
        s,
        "Live Demo — 60 Seconds",
        "Problem → solution → impact in one linear flow",
        criterion="JUDGING CRITERION 5 — Presentation & pitch",
    )
    steps = [
        ("0:00", "Zone map all GREEN — normal underground shift"),
        ("0:15", "Zone 3 methane CLIMBING → YELLOW + rising arrow (the trend)"),
        ("0:30", "Fan ramps — ventilation actively mitigating"),
        ("0:40", "Crosses 1.5% → RED · banner · voice reads Reg 854 procedure"),
        ("0:50", "Turn WIFI OFF · run demo again · still works (offline proof)"),
        ("1:00", "Zone clears → \"Zone 3 stabilized. Everything under control.\""),
    ]
    y = 1.85
    for time, desc in steps:
        add_textbox(s, 0.9, y, 1.2, 0.35, time, size=14, bold=True, color=ACCENT)
        add_textbox(s, 2.2, y, 9.5, 0.35, desc, size=16)
        y += 0.52
    add_textbox(s, 0.9, 6.0, 11, 0.5, "Backup: pre-recorded demo video if live fails", size=12, color=MUTED)

    # ── CRITERION 4: Impact & feasibility ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(
        s,
        "Impact & Feasibility",
        "Meaningful value with a clear path to deployment",
        criterion="JUDGING CRITERION 4 — Impact & feasibility",
    )
    add_bullets(s, 0.9, 1.85, 5.8, 4.0, [
        "Early warning → orderly withdrawal, not emergency scramble",
        "Lives protected — cited Reg 854 procedures reduce human error",
        "Works when comms fail — offline edge, no central server",
        "Power savings — fans idle until gas builds, then ramp",
        "Auditable — every alert cites O. Reg 854 section + text",
    ], size=16)
    add_textbox(s, 7.0, 1.85, 5.5, 0.4, "Next steps to deploy", size=16, bold=True, color=GREEN)
    add_bullets(s, 7.0, 2.35, 5.5, 3.5, [
        "Phase 1: Deploy edge box per zone on existing mesh",
        "Phase 2: Connect real methane / CO sensor hardware",
        "Phase 3: Integrate cap-lamp TTS for hands-free alerts",
        "Phase 4: Surface sync for incident logs & compliance",
    ], size=15)

    # ── CRITERION 6: Responsible design & limitations ──────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(
        s,
        "Responsible Design & Limitations",
        "We know what we built — and what we did not",
        criterion="JUDGING CRITERION 6 — Responsible design & awareness of limitations",
    )
    add_textbox(s, 0.9, 1.85, 5.5, 0.4, "Safety & operational awareness", size=16, bold=True, color=GREEN)
    add_bullets(s, 0.9, 2.3, 5.5, 3.5, [
        "Never enrich oxygen underground — only increase airflow",
        "Demo thresholds (1.0% / 1.5%) are illustrative; real legal limits differ",
        "Template answers, not a cloud LLM — auditable but not open-ended",
        "No personal data collected — sensor readings only, local storage",
    ], size=15)
    add_textbox(s, 7.0, 1.85, 5.5, 0.4, "Honest limitations", size=16, bold=True, color=ACCENT)
    add_bullets(s, 7.0, 2.3, 5.5, 3.5, [
        "Today: simulated zones on one laptop — not real hardware mesh",
        "No real Bluetooth/radio mesh shipped in 3 hours",
        "Web Speech API varies by browser — on-screen text is fallback",
        "Reg 854 retrieval is keyword-based, not full legal AI",
    ], size=15)

    # ── CRITERION 7: Originality & effort ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(
        s,
        "Originality & Effort",
        "Thoughtful problem-solving with clear implementation work",
        criterion="JUDGING CRITERION 7 — Originality & effort",
    )
    add_bullets(s, 0.9, 1.85, 11.2, 4.5, [
        "Original: trend-based early warning + automated ventilation loop + cited voice alerts — combined offline",
        "Not just a dashboard — the system ACTS (fan ramp) and SPEAKS (Reg 854 procedure)",
        "Closed-loop simulator: gas inrush → danger → venting → recovery → repeat (demo-ready anytime)",
        "Full vertical slice in 3 hours: MQTT · FastAPI · trend engine · Reg 854 · UI · voice · pitch deck",
        "Built for the real constraint (no internet underground) — not retrofitted from a cloud app",
    ], size=17)
    add_textbox(s, 0.9, 5.5, 11.5, 0.8,
        "Team split: Backend/Edge (B) · Frontend (F) · Voice + Pitch (V) — minimal overlap, maximum output",
        size=14, color=MUTED, align=PP_ALIGN.CENTER,
    )

    # ── JUDGING CRITERIA SCORECARD (summary for judges) ────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(s, "How We Meet Every Criterion", "Quick reference for judges")
    criteria = [
        ("1. Problem & track alignment", "Mining-track gas safety; also Health + Sustainability"),
        ("2. Solution design", "Edge offline system for real supervisors and crews"),
        ("3. Working prototype", "Live API, 4 zones, trend + ventilation + Reg 854 — demo now"),
        ("4. Impact & feasibility", "Saves lives, saves power; 4-phase deployment path"),
        ("5. Presentation & pitch", "60s demo: green → trend → red → wifi off → all clear"),
        ("6. Responsible design", "Safety-aware; honest about demo vs production scope"),
        ("7. Originality & effort", "Trend + act + speak offline; full stack in 3 hours"),
    ]
    y = 1.75
    for label, answer in criteria:
        add_textbox(s, 0.9, y, 4.5, 0.35, label, size=13, bold=True, color=ACCENT)
        add_textbox(s, 5.5, y, 6.8, 0.35, answer, size=13)
        y += 0.52

    # ── BUILD TODAY vs PRODUCTION ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    slide_header(s, "What We Built vs. Production Vision", "Transparent scope")
    add_textbox(s, 0.9, 1.7, 5.5, 0.4, "TODAY (hackathon build)", size=18, bold=True, color=GREEN)
    add_bullets(s, 0.9, 2.15, 5.5, 3.5, [
        "Simulated zones on one laptop",
        "Local MQTT or internal timer fallback",
        "Full API + trend + ventilation + Reg 854",
        "Offline demo with wifi disabled",
    ], size=15)
    add_textbox(s, 7.0, 1.7, 5.5, 0.4, "PRODUCTION (deployment)", size=18, bold=True, color=ACCENT)
    add_bullets(s, 7.0, 2.15, 5.5, 3.5, [
        "Real sensor node per zone on mesh network",
        "Wireless self-healing transport (no leaky feeder)",
        "Edge box underground, sync at surface",
        "On-device TTS in cap lamp / rugged tablet",
    ], size=15)

    # ── CLOSE ──────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.9, 2.0, 11.5, 1.0, "MineGuard Edge", size=44, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        s, 0.9, 3.1, 11.5, 1.0,
        "Get crews out before it's an emergency — offline, cited, and spoken.",
        size=20, color=ACCENT, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, 0.9, 4.2, 11.5, 0.6,
        "Mining track  ·  Worker safety  ·  Smart ventilation",
        size=16, color=GREEN, align=PP_ALIGN.CENTER,
    )
    add_textbox(s, 0.9, 5.2, 11.5, 0.6, "Questions?", size=28, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        s, 0.9, 6.0, 11.5, 0.5,
        "github.com/RavalDH/Cursor_Hack  ·  localhost:8000/docs",
        size=13, color=MUTED, align=PP_ALIGN.CENTER,
    )

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
